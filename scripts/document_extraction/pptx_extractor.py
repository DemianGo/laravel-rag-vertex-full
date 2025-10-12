#!/usr/bin/env python3
"""
PPTX Extractor Wrapper
Directly extracts PPTX files, works with temp files without extensions.
"""

import sys
import os
import shutil

def extract_pptx(file_path: str) -> str:
    """Extract text from PPTX"""
    temp_link = None
    try:
        if not os.path.exists(file_path):
            return ''
        
        # python-pptx requires .pptx extension
        if not file_path.endswith('.pptx') and not file_path.endswith('.ppt'):
            temp_link = file_path + '.pptx'
            shutil.copy2(file_path, temp_link)
            file_to_read = temp_link
        else:
            file_to_read = file_path
        
        from pptx import Presentation
        
        prs = Presentation(file_to_read)
        
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"=== Slide {slide_num} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
        
        return '\n'.join(text_parts)
    except Exception as e:
        # Silent fail - return empty
        return ''
    finally:
        # Clean up temp file
        if temp_link and os.path.exists(temp_link):
            try:
                os.remove(temp_link)
            except:
                pass

if __name__ == "__main__":
    if len(sys.argv) < 2:
        sys.exit(1)
    
    file_path = sys.argv[1]
    text = extract_pptx(file_path)
    print(text)

