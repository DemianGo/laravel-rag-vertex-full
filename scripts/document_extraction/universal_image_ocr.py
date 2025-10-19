#!/usr/bin/env python3
"""
Universal Image OCR Processor
Extracts and processes images from all document formats (DOCX, PPTX, XLSX, HTML, XML, RTF)
using Google Vision as default OCR engine.
"""

import os
import sys
import json
import tempfile
import shutil
from pathlib import Path
from typing import Dict, List, Any, Optional
import logging

# Import document processing libraries
try:
    from docx import Document
    from docx.parts.image import ImagePart
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    HTML_AVAILABLE = True
except ImportError:
    HTML_AVAILABLE = False

# Import our OCR processor
try:
    from advanced_ocr_processor import AdvancedOCRProcessor
    ADVANCED_OCR_AVAILABLE = True
except ImportError:
    ADVANCED_OCR_AVAILABLE = False

# Import image processing
try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


class UniversalImageOCR:
    """
    Universal Image OCR Processor that extracts and processes images from all document formats.
    Uses Google Vision as default OCR engine with Tesseract fallback.
    """
    
    def __init__(self, use_google_vision: bool = True):
        """Initialize the universal OCR processor."""
        self.use_google_vision = use_google_vision and ADVANCED_OCR_AVAILABLE
        self.temp_dir = None
        self.advanced_ocr = None
        
        if self.use_google_vision:
            try:
                self.advanced_ocr = AdvancedOCRProcessor(use_google_vision=True)
            except Exception as e:
                logging.warning(f"Failed to initialize AdvancedOCRProcessor: {e}")
                self.use_google_vision = False
        
        # Create temporary directory for image extraction
        self.temp_dir = tempfile.mkdtemp(prefix='ocr_images_')
    
    def extract_and_process_images(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """
        Extract and process images from document.
        
        Args:
            file_path: Path to the document file
            file_type: Type of document (docx, pptx, xlsx, html, xml, rtf)
            
        Returns:
            Dict with success status, OCR text, and metadata
        """
        try:
            if not self.use_google_vision or not self.advanced_ocr:
                return {
                    "success": False,
                    "error": "Google Vision OCR not available",
                    "ocr_text": "",
                    "images_processed": 0
                }
            
            # Extract images based on file type
            image_paths = []
            
            if file_type == 'docx' and DOCX_AVAILABLE:
                image_paths = self._extract_docx_images(file_path)
            elif file_type == 'pptx' and PPTX_AVAILABLE:
                image_paths = self._extract_pptx_images(file_path)
            elif file_type == 'xlsx' and XLSX_AVAILABLE:
                image_paths = self._extract_xlsx_images(file_path)
            elif file_type in ['html', 'xml'] and HTML_AVAILABLE:
                image_paths = self._extract_html_images(file_path)
            elif file_type == 'rtf':
                image_paths = self._extract_rtf_images(file_path)
            else:
                return {
                    "success": False,
                    "error": f"Unsupported file type for image extraction: {file_type}",
                    "ocr_text": "",
                    "images_processed": 0
                }
            
            if not image_paths:
                return {
                    "success": True,
                    "ocr_text": "",
                    "images_processed": 0,
                    "message": "No images found in document"
                }
            
            # Process images with Google Vision
            ocr_results = []
            processed_count = 0
            
            for image_path in image_paths:
                try:
                    result = self.advanced_ocr.process_image(str(image_path))
                    if result.get('success') and result.get('text'):
                        ocr_results.append(result['text'])
                        processed_count += 1
                except Exception as e:
                    logging.warning(f"Failed to process image {image_path}: {e}")
                    continue
            
            # Combine all OCR text
            combined_ocr_text = '\n\n'.join(ocr_results)
            
            return {
                "success": True,
                "ocr_text": combined_ocr_text,
                "images_processed": processed_count,
                "total_images": len(image_paths),
                "temp_dir": self.temp_dir
            }
            
        except Exception as e:
            logging.error(f"Error in extract_and_process_images: {e}")
            return {
                "success": False,
                "error": f"Image extraction failed: {str(e)}",
                "ocr_text": "",
                "images_processed": 0
            }
        finally:
            # Cleanup temporary files
            self.cleanup()
    
    def _extract_docx_images(self, file_path: str) -> List[Path]:
        """Extract images from DOCX file."""
        image_paths = []
        
        try:
            doc = Document(file_path)
            
            # Extract images from relationships
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        # Get image data
                        image_data = rel.target_part.blob
                        
                        # Determine image extension
                        ext = rel.target_ref.split('.')[-1].lower()
                        if ext not in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp']:
                            ext = 'png'  # Default fallback
                        
                        # Save image to temp directory
                        image_filename = f"docx_image_{len(image_paths)}_{rel.rId}.{ext}"
                        image_path = Path(self.temp_dir) / image_filename
                        
                        with open(image_path, 'wb') as img_file:
                            img_file.write(image_data)
                        
                        image_paths.append(image_path)
                        
                    except Exception as e:
                        logging.warning(f"Failed to extract DOCX image {rel.rId}: {e}")
                        continue
            
            return image_paths
            
        except Exception as e:
            logging.error(f"Error extracting DOCX images: {e}")
            return []
    
    def _extract_pptx_images(self, file_path: str) -> List[Path]:
        """Extract images from PPTX file."""
        image_paths = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, 'image'):
                        try:
                            # Get image data
                            image_data = shape.image.blob
                            
                            # Determine image extension
                            ext = shape.image.ext
                            if not ext or ext not in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp']:
                                ext = 'png'  # Default fallback
                            
                            # Save image to temp directory
                            image_filename = f"pptx_slide_{slide_idx}_image_{len(image_paths)}.{ext}"
                            image_path = Path(self.temp_dir) / image_filename
                            
                            with open(image_path, 'wb') as img_file:
                                img_file.write(image_data)
                            
                            image_paths.append(image_path)
                            
                        except Exception as e:
                            logging.warning(f"Failed to extract PPTX image from slide {slide_idx}: {e}")
                            continue
            
            return image_paths
            
        except Exception as e:
            logging.error(f"Error extracting PPTX images: {e}")
            return []
    
    def _extract_pptx_slide_images(self, file_path: str, slide_number: int) -> List[Path]:
        """Extract images from a specific slide in PPTX file."""
        image_paths = []
        
        try:
            prs = Presentation(file_path)
            
            if slide_number < len(prs.slides):
                slide = prs.slides[slide_number]
                
                for shape in slide.shapes:
                    if hasattr(shape, 'image'):
                        try:
                            # Get image data
                            image_data = shape.image.blob
                            
                            # Determine image extension
                            ext = shape.image.ext
                            if not ext or ext not in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp']:
                                ext = 'png'  # Default fallback
                            
                            # Save image to temp directory
                            image_filename = f"pptx_slide_{slide_number}_image_{len(image_paths)}.{ext}"
                            image_path = Path(self.temp_dir) / image_filename
                            
                            with open(image_path, 'wb') as img_file:
                                img_file.write(image_data)
                            
                            image_paths.append(image_path)
                            
                        except Exception as e:
                            logging.warning(f"Failed to extract PPTX image from slide {slide_number}: {e}")
                            continue
            
            return image_paths
            
        except Exception as e:
            logging.error(f"Error extracting PPTX slide images: {e}")
            return []
    
    def _extract_xlsx_images(self, file_path: str) -> List[Path]:
        """Extract images from XLSX file."""
        image_paths = []
        
        try:
            workbook = load_workbook(file_path)
            
            # XLSX images are stored in drawing parts
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Check for images in sheet
                if hasattr(sheet, '_images'):
                    for img_idx, image in enumerate(sheet._images):
                        try:
                            # Save image to temp directory
                            image_filename = f"xlsx_{sheet_name}_image_{img_idx}.png"
                            image_path = Path(self.temp_dir) / image_filename
                            
                            # Convert and save image
                            if PIL_AVAILABLE:
                                img = Image.frombytes('RGBA', image.size, image.data)
                                img.save(image_path)
                                image_paths.append(image_path)
                            
                        except Exception as e:
                            logging.warning(f"Failed to extract XLSX image from sheet {sheet_name}: {e}")
                            continue
            
            return image_paths
            
        except Exception as e:
            logging.error(f"Error extracting XLSX images: {e}")
            return []
    
    def _extract_html_images(self, file_path: str) -> List[Path]:
        """Extract images from HTML/XML file."""
        image_paths = []
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            
            # Find all img tags
            for img_tag in soup.find_all('img'):
                src = img_tag.get('src')
                if not src:
                    continue
                
                # Handle base64 images
                if src.startswith('data:image'):
                    try:
                        # Parse base64 image
                        header, data = src.split(',', 1)
                        import base64
                        
                        # Determine extension from header
                        if 'jpeg' in header or 'jpg' in header:
                            ext = 'jpg'
                        elif 'png' in header:
                            ext = 'png'
                        elif 'gif' in header:
                            ext = 'gif'
                        else:
                            ext = 'png'  # Default
                        
                        # Decode and save
                        image_data = base64.b64decode(data)
                        image_filename = f"html_image_{len(image_paths)}.{ext}"
                        image_path = Path(self.temp_dir) / image_filename
                        
                        with open(image_path, 'wb') as img_file:
                            img_file.write(image_data)
                        
                        image_paths.append(image_path)
                        
                    except Exception as e:
                        logging.warning(f"Failed to extract base64 image: {e}")
                        continue
                
                # Handle local file references (if file exists)
                elif src.startswith('/') or not src.startswith('http'):
                    local_path = Path(file_path).parent / src
                    if local_path.exists():
                        image_filename = f"html_image_{len(image_paths)}{local_path.suffix}"
                        image_path = Path(self.temp_dir) / image_filename
                        
                        try:
                            shutil.copy2(local_path, image_path)
                            image_paths.append(image_path)
                        except Exception as e:
                            logging.warning(f"Failed to copy local image {src}: {e}")
                            continue
            
            return image_paths
            
        except Exception as e:
            logging.error(f"Error extracting HTML images: {e}")
            return []
    
    def _extract_rtf_images(self, file_path: str) -> List[Path]:
        """Extract images from RTF file."""
        image_paths = []
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # RTF images are encoded in hexadecimal format
            # This is a simplified implementation - RTF parsing is complex
            import re
            
            # Look for RTF image objects
            image_pattern = r'{\\pict[^}]*}([^}]+)'
            matches = re.findall(image_pattern, content)
            
            for idx, match in enumerate(matches):
                try:
                    # Clean hex data
                    hex_data = re.sub(r'[^0-9a-fA-F]', '', match)
                    if len(hex_data) % 2 != 0:
                        continue
                    
                    # Convert hex to bytes
                    image_data = bytes.fromhex(hex_data)
                    
                    # Save image (assuming PNG format for RTF)
                    image_filename = f"rtf_image_{idx}.png"
                    image_path = Path(self.temp_dir) / image_filename
                    
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_data)
                    
                    image_paths.append(image_path)
                    
                except Exception as e:
                    logging.warning(f"Failed to extract RTF image {idx}: {e}")
                    continue
            
            return image_paths
            
        except Exception as e:
            logging.error(f"Error extracting RTF images: {e}")
            return []
    
    def cleanup(self):
        """Clean up temporary files."""
        if self.temp_dir and Path(self.temp_dir).exists():
            try:
                shutil.rmtree(self.temp_dir)
            except Exception as e:
                logging.warning(f"Failed to cleanup temp directory {self.temp_dir}: {e}")


def main():
    """CLI interface for UniversalImageOCR."""
    if len(sys.argv) < 3:
        print(json.dumps({
            "success": False,
            "error": "Usage: universal_image_ocr.py <file_path> <file_type>"
        }))
        sys.exit(1)
    
    file_path = sys.argv[1]
    file_type = sys.argv[2]
    
    if not Path(file_path).exists():
        print(json.dumps({
            "success": False,
            "error": f"File not found: {file_path}"
        }))
        sys.exit(1)
    
    processor = UniversalImageOCR(use_google_vision=True)
    result = processor.extract_and_process_images(file_path, file_type)
    
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
