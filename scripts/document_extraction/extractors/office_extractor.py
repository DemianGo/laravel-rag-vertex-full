#!/usr/bin/env python3
"""
Office Documents Extractor

Handles DOCX, XLSX, PPTX files using appropriate libraries.
"""

import os
from pathlib import Path


class OfficeExtractor:
    def __init__(self):
        pass

    def extract(self, file_path, file_type):
        """Extract content based on office file type"""
        try:
            if file_type == 'docx':
                return self._extract_docx(file_path)
            elif file_type == 'xlsx':
                return self._extract_xlsx(file_path)
            elif file_type == 'pptx':
                return self._extract_pptx(file_path)
            else:
                return self._error_response(f"Unsupported office type: {file_type}")

        except Exception as e:
            return self._error_response(f"Office extraction failed: {str(e)}")

    def _extract_docx(self, file_path):
        """Extract text from DOCX files"""
        try:
            from docx import Document

            doc = Document(file_path)
            extracted_text = []
            total_elements = 0
            extracted_elements = 0
            failed_elements = 0

            # Extract paragraphs
            for para in doc.paragraphs:
                total_elements += 1
                try:
                    text = para.text.strip()
                    if text:
                        extracted_text.append(text)
                        extracted_elements += 1
                except Exception:
                    failed_elements += 1

            # Extract tables
            for table in doc.tables:
                total_elements += 1
                try:
                    table_text = []
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_text.append(cell_text)
                        if row_text:
                            table_text.append(" | ".join(row_text))

                    if table_text:
                        extracted_text.append("\n".join(table_text))
                        extracted_elements += 1
                except Exception:
                    failed_elements += 1

            # Extract headers and footers
            for section in doc.sections:
                total_elements += 1
                try:
                    header_text = []
                    footer_text = []

                    # Header
                    if section.header:
                        for para in section.header.paragraphs:
                            text = para.text.strip()
                            if text:
                                header_text.append(text)

                    # Footer
                    if section.footer:
                        for para in section.footer.paragraphs:
                            text = para.text.strip()
                            if text:
                                footer_text.append(text)

                    if header_text:
                        extracted_text.append("HEADER: " + " ".join(header_text))
                    if footer_text:
                        extracted_text.append("FOOTER: " + " ".join(footer_text))

                    if header_text or footer_text:
                        extracted_elements += 1

                except Exception:
                    failed_elements += 1

            combined_text = "\n\n".join(extracted_text)

            return {
                'status': 'success',
                'extracted_text': combined_text,
                'metadata': {
                    'document_type': 'docx',
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables),
                    'sections': len(doc.sections),
                    'total_characters': len(combined_text)
                },
                'metrics': {
                    'total_elements': total_elements,
                    'extracted_elements': extracted_elements,
                    'failed_elements': failed_elements,
                    'extraction_percentage': (extracted_elements / total_elements * 100) if total_elements > 0 else 0
                }
            }

        except ImportError:
            return self._error_response("python-docx library not installed. Run: pip install python-docx")
        except Exception as e:
            return self._error_response(f"DOCX extraction error: {str(e)}")

    def _extract_xlsx(self, file_path):
        """Extract data from XLSX files"""
        try:
            from openpyxl import load_workbook

            workbook = load_workbook(file_path, data_only=True)
            extracted_text = []
            total_elements = 0
            extracted_elements = 0
            failed_elements = 0

            sheet_data = {}

            for sheet_name in workbook.sheetnames:
                total_elements += 1
                try:
                    worksheet = workbook[sheet_name]
                    sheet_text = []
                    cell_count = 0

                    for row in worksheet.iter_rows(values_only=True):
                        row_text = []
                        for cell_value in row:
                            if cell_value is not None:
                                row_text.append(str(cell_value))
                                cell_count += 1

                        if row_text:
                            sheet_text.append(" | ".join(row_text))

                    if sheet_text:
                        sheet_content = f"SHEET: {sheet_name}\n" + "\n".join(sheet_text)
                        extracted_text.append(sheet_content)
                        sheet_data[sheet_name] = {
                            'rows': len(sheet_text),
                            'cells_with_data': cell_count
                        }
                        extracted_elements += 1

                except Exception:
                    failed_elements += 1

            combined_text = "\n\n".join(extracted_text)

            return {
                'status': 'success',
                'extracted_text': combined_text,
                'metadata': {
                    'document_type': 'xlsx',
                    'total_sheets': len(workbook.sheetnames),
                    'sheet_names': list(workbook.sheetnames),
                    'sheet_data': sheet_data,
                    'total_characters': len(combined_text)
                },
                'metrics': {
                    'total_elements': total_elements,
                    'extracted_elements': extracted_elements,
                    'failed_elements': failed_elements,
                    'extraction_percentage': (extracted_elements / total_elements * 100) if total_elements > 0 else 0
                }
            }

        except ImportError:
            return self._error_response("openpyxl library not installed. Run: pip install openpyxl")
        except Exception as e:
            return self._error_response(f"XLSX extraction error: {str(e)}")

    def _extract_pptx(self, file_path):
        """Extract text from PPTX files"""
        try:
            from pptx import Presentation

            presentation = Presentation(file_path)
            extracted_text = []
            total_elements = 0
            extracted_elements = 0
            failed_elements = 0

            slide_data = []

            for slide_num, slide in enumerate(presentation.slides, 1):
                total_elements += 1
                try:
                    slide_text = []

                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())

                    # Extract notes
                    if slide.notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_text.append(f"NOTES: {notes_text}")

                    if slide_text:
                        slide_content = f"SLIDE {slide_num}:\n" + "\n".join(slide_text)
                        extracted_text.append(slide_content)
                        slide_data.append({
                            'slide_number': slide_num,
                            'text_shapes': len([s for s in slide.shapes if hasattr(s, "text") and s.text.strip()]),
                            'has_notes': bool(slide.notes_slide and slide.notes_slide.notes_text_frame and slide.notes_slide.notes_text_frame.text.strip())
                        })
                        extracted_elements += 1

                except Exception:
                    failed_elements += 1

            combined_text = "\n\n".join(extracted_text)

            return {
                'status': 'success',
                'extracted_text': combined_text,
                'metadata': {
                    'document_type': 'pptx',
                    'total_slides': len(presentation.slides),
                    'slides_with_content': len(slide_data),
                    'slide_data': slide_data,
                    'total_characters': len(combined_text)
                },
                'metrics': {
                    'total_elements': total_elements,
                    'extracted_elements': extracted_elements,
                    'failed_elements': failed_elements,
                    'extraction_percentage': (extracted_elements / total_elements * 100) if total_elements > 0 else 0
                }
            }

        except ImportError:
            return self._error_response("python-pptx library not installed. Run: pip install python-pptx")
        except Exception as e:
            return self._error_response(f"PPTX extraction error: {str(e)}")

    def _error_response(self, error_msg):
        """Standardized error response"""
        return {
            'status': 'error',
            'error': error_msg,
            'extracted_text': '',
            'metadata': {},
            'metrics': {
                'total_elements': 0,
                'extracted_elements': 0,
                'failed_elements': 0,
                'extraction_percentage': 0.0
            }
        }