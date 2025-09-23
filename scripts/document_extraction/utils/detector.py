"""
Advanced document type detection based on content analysis.
"""

import os
import mimetypes
from pathlib import Path
from typing import Dict, List, Any, Optional

try:
    import magic
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False


def detect_document_type(file_path: str) -> Dict[str, Any]:
    """
    Advanced document type detection using magic bytes and content analysis.

    Args:
        file_path: Path to the file to analyze

    Returns:
        Dict with primary_type, sub_types, and confidence score
    """
    if not os.path.exists(file_path):
        return {
            "primary_type": "unknown",
            "sub_types": [],
            "confidence": 0.0,
            "error": "File not found"
        }

    try:
        # Get file extension
        file_ext = Path(file_path).suffix.lower()

        # Use magic library if available
        mime_type = None
        magic_type = None

        if MAGIC_AVAILABLE:
            try:
                mime = magic.Magic(mime=True)
                mime_type = mime.from_file(file_path)

                magic_desc = magic.Magic()
                magic_type = magic_desc.from_file(file_path)
            except Exception:
                pass

        # Fallback to Python's mimetypes
        if not mime_type:
            mime_type, _ = mimetypes.guess_type(file_path)

        # Analyze content for sub-types
        sub_types = []
        confidence = 0.8

        # Determine primary type
        primary_type = _determine_primary_type(file_ext, mime_type, magic_type)

        # Analyze content for hybrid detection
        if primary_type == "pdf":
            sub_types, conf_modifier = _analyze_pdf_content(file_path)
            confidence += conf_modifier
        elif primary_type in ["docx", "xlsx", "pptx"]:
            sub_types, conf_modifier = _analyze_office_content(file_path, primary_type)
            confidence += conf_modifier
        elif primary_type in ["txt", "csv", "rtf"]:
            sub_types, conf_modifier = _analyze_text_content(file_path)
            confidence += conf_modifier
        elif primary_type in ["html", "xml"]:
            sub_types, conf_modifier = _analyze_web_content(file_path)
            confidence += conf_modifier

        # Ensure confidence is within bounds
        confidence = min(1.0, max(0.1, confidence))

        return {
            "primary_type": primary_type,
            "sub_types": sub_types,
            "confidence": confidence,
            "mime_type": mime_type,
            "magic_type": magic_type,
            "file_extension": file_ext
        }

    except Exception as e:
        return {
            "primary_type": "unknown",
            "sub_types": [],
            "confidence": 0.0,
            "error": str(e)
        }


def _determine_primary_type(file_ext: str, mime_type: Optional[str], magic_type: Optional[str]) -> str:
    """Determine the primary document type."""

    # PDF detection
    if (file_ext == ".pdf" or
        (mime_type and "pdf" in mime_type.lower()) or
        (magic_type and "PDF" in magic_type)):
        return "pdf"

    # Office document detection
    office_extensions = {
        ".docx": "docx",
        ".doc": "docx",
        ".xlsx": "xlsx",
        ".xls": "xlsx",
        ".pptx": "pptx",
        ".ppt": "pptx"
    }

    if file_ext in office_extensions:
        return office_extensions[file_ext]

    if mime_type:
        if "wordprocessingml" in mime_type or "msword" in mime_type:
            return "docx"
        elif "spreadsheetml" in mime_type or "excel" in mime_type:
            return "xlsx"
        elif "presentationml" in mime_type or "powerpoint" in mime_type:
            return "pptx"

    # Web document detection (prioritize before text)
    web_extensions = {".html": "html", ".htm": "html", ".xml": "xml", ".xhtml": "html"}
    if file_ext in web_extensions:
        return web_extensions[file_ext]

    if mime_type:
        if "html" in mime_type:
            return "html"
        elif "xml" in mime_type:
            return "xml"

    # Text file detection
    text_extensions = [".txt", ".csv", ".rtf", ".tsv"]
    if file_ext in text_extensions:
        return file_ext[1:]  # Remove dot

    if mime_type and mime_type.startswith("text/"):
        if "csv" in mime_type:
            return "csv"
        elif "rtf" in mime_type:
            return "rtf"
        else:
            return "txt"

    return "unknown"


def _analyze_pdf_content(file_path: str) -> tuple[List[str], float]:
    """Analyze PDF content to detect sub-types."""
    sub_types = []
    confidence_modifier = 0.0

    try:
        from PyPDF2 import PdfReader

        with open(file_path, 'rb') as file:
            reader = PdfReader(file)

            # Check for text content
            has_text = False
            text_pages = 0

            for i, page in enumerate(reader.pages[:5]):  # Check first 5 pages
                text = page.extract_text().strip()
                if text and len(text) > 50:
                    has_text = True
                    text_pages += 1

            if has_text:
                sub_types.append("text")
                confidence_modifier += 0.1

            # Check for images (indirect detection)
            total_pages = len(reader.pages)
            if text_pages < total_pages * 0.5:  # Less than 50% pages have substantial text
                sub_types.append("images")
                confidence_modifier += 0.05

            # Check for forms
            if reader.metadata and "/AcroForm" in str(reader.metadata):
                sub_types.append("forms")
                confidence_modifier += 0.05

    except Exception:
        # Fallback - assume it has text
        sub_types.append("text")

    return sub_types, confidence_modifier


def _analyze_office_content(file_path: str, doc_type: str) -> tuple[List[str], float]:
    """Analyze Office document content."""
    sub_types = []
    confidence_modifier = 0.0

    try:
        if doc_type == "docx":
            from docx import Document
            doc = Document(file_path)

            # Check for tables
            if doc.tables:
                sub_types.append("tables")
                confidence_modifier += 0.1

            # Check for images
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    sub_types.append("images")
                    confidence_modifier += 0.05
                    break

            # Always has text
            sub_types.append("text")

        elif doc_type == "xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True)

            # Check for multiple sheets
            if len(wb.sheetnames) > 1:
                sub_types.append("multi_sheet")
                confidence_modifier += 0.05

            # Check for formulas
            ws = wb.active
            has_formulas = False
            for row in ws.iter_rows(max_row=100, max_col=20):  # Sample check
                for cell in row:
                    if cell.value and str(cell.value).startswith('='):
                        has_formulas = True
                        break
                if has_formulas:
                    break

            if has_formulas:
                sub_types.append("formulas")
                confidence_modifier += 0.1

            sub_types.append("data")

        elif doc_type == "pptx":
            from pptx import Presentation
            pres = Presentation(file_path)

            # Check slide count
            if len(pres.slides) > 10:
                sub_types.append("long_presentation")
                confidence_modifier += 0.05

            # Check for images/media
            has_media = False
            for slide in pres.slides[:5]:  # Check first 5 slides
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        has_media = True
                        break
                if has_media:
                    break

            if has_media:
                sub_types.append("media")
                confidence_modifier += 0.1

            sub_types.append("slides")

    except Exception:
        # Fallback
        if doc_type == "docx":
            sub_types.append("text")
        elif doc_type == "xlsx":
            sub_types.append("data")
        elif doc_type == "pptx":
            sub_types.append("slides")

    return sub_types, confidence_modifier


def _analyze_text_content(file_path: str) -> tuple[List[str], float]:
    """Analyze text file content."""
    sub_types = []
    confidence_modifier = 0.0

    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        content = None

        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read(10000)  # Read first 10KB
                break
            except UnicodeDecodeError:
                continue

        if content:
            # Check for structured data
            if '\t' in content or ',' in content:
                lines = content.split('\n')[:10]
                if len(lines) > 1:
                    # Check if it looks like CSV/TSV
                    delimiters = [',', '\t', ';', '|']
                    for delimiter in delimiters:
                        if all(delimiter in line for line in lines if line.strip()):
                            sub_types.append("structured")
                            confidence_modifier += 0.1
                            break

            # Check for markup/code patterns
            if any(marker in content.lower() for marker in ['<html', '<xml', '<?xml', '<body']):
                sub_types.append("markup")
                confidence_modifier += 0.1
            elif any(keyword in content for keyword in ['def ', 'function ', 'class ', 'import ']):
                sub_types.append("code")
                confidence_modifier += 0.1

            sub_types.append("plain_text")

    except Exception:
        sub_types.append("plain_text")

    return sub_types, confidence_modifier


def _analyze_web_content(file_path: str) -> tuple[List[str], float]:
    """Analyze web document content."""
    sub_types = []
    confidence_modifier = 0.0

    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read(10000).lower()

        # Check for HTML features
        if any(tag in content for tag in ['<table', '<form', '<script']):
            if '<table' in content:
                sub_types.append("tables")
            if '<form' in content:
                sub_types.append("forms")
            if '<script' in content:
                sub_types.append("interactive")
            confidence_modifier += 0.1

        # Check for XML features
        if content.strip().startswith('<?xml'):
            sub_types.append("structured_data")
            confidence_modifier += 0.1

        if not sub_types:
            sub_types.append("markup")

    except Exception:
        sub_types.append("markup")

    return sub_types, confidence_modifier