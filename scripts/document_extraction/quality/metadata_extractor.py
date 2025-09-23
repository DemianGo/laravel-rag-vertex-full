"""
Metadata extraction for different document types.
"""

import os
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, Optional


class MetadataExtractor:
    """Extract metadata from various document types."""

    def extract_metadata(self, file_path: str, doc_type: str) -> Dict[str, Any]:
        """
        Extract metadata based on document type.

        Args:
            file_path: Path to the document
            doc_type: Type of document (pdf, docx, xlsx, etc.)

        Returns:
            Dictionary with extracted metadata
        """
        base_metadata = self._get_base_metadata(file_path)

        if doc_type == "pdf":
            specific_metadata = self.extract_pdf_metadata(file_path)
        elif doc_type in ["docx", "xlsx", "pptx"]:
            specific_metadata = self.extract_office_metadata(file_path, doc_type)
        elif doc_type in ["txt", "csv", "rtf"]:
            specific_metadata = self.extract_text_metadata(file_path)
        elif doc_type in ["html", "xml"]:
            specific_metadata = self.extract_web_metadata(file_path)
        else:
            specific_metadata = {}

        # Merge base and specific metadata
        metadata = {**base_metadata, **specific_metadata}
        return metadata

    def _get_base_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract basic file system metadata."""
        try:
            stat = os.stat(file_path)
            path_obj = Path(file_path)

            return {
                "filename": path_obj.name,
                "file_size": stat.st_size,
                "file_extension": path_obj.suffix.lower(),
                "created_at": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "modified_at": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "accessed_at": datetime.fromtimestamp(stat.st_atime).isoformat()
            }
        except Exception as e:
            return {"error": f"Failed to extract base metadata: {str(e)}"}

    def extract_pdf_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract PDF-specific metadata."""
        try:
            from PyPDF2 import PdfReader

            metadata = {}

            with open(file_path, 'rb') as file:
                reader = PdfReader(file)

                # Basic document info
                metadata["page_count"] = len(reader.pages)

                # PDF metadata
                if reader.metadata:
                    pdf_meta = reader.metadata
                    metadata["title"] = pdf_meta.get("/Title", "").strip()
                    metadata["author"] = pdf_meta.get("/Author", "").strip()
                    metadata["subject"] = pdf_meta.get("/Subject", "").strip()
                    metadata["creator"] = pdf_meta.get("/Creator", "").strip()
                    metadata["producer"] = pdf_meta.get("/Producer", "").strip()

                    # Creation and modification dates
                    if "/CreationDate" in pdf_meta:
                        metadata["creation_date"] = str(pdf_meta["/CreationDate"])
                    if "/ModDate" in pdf_meta:
                        metadata["modification_date"] = str(pdf_meta["/ModDate"])

                # Analyze content structure
                text_pages = 0
                total_chars = 0

                for page in reader.pages[:5]:  # Sample first 5 pages
                    try:
                        text = page.extract_text()
                        if text and len(text.strip()) > 50:
                            text_pages += 1
                            total_chars += len(text)
                    except Exception:
                        continue

                metadata["text_pages_sample"] = text_pages
                metadata["avg_chars_per_page"] = total_chars // max(text_pages, 1)

                # Check for security features
                metadata["encrypted"] = reader.is_encrypted
                if hasattr(reader, 'security_handler'):
                    metadata["security_handler"] = str(type(reader.security_handler).__name__)

            return metadata

        except Exception as e:
            return {"error": f"Failed to extract PDF metadata: {str(e)}"}

    def extract_office_metadata(self, file_path: str, doc_type: str) -> Dict[str, Any]:
        """Extract metadata from Office documents."""
        try:
            if doc_type == "docx":
                return self._extract_docx_metadata(file_path)
            elif doc_type == "xlsx":
                return self._extract_xlsx_metadata(file_path)
            elif doc_type == "pptx":
                return self._extract_pptx_metadata(file_path)
            else:
                return {"error": f"Unsupported Office document type: {doc_type}"}

        except Exception as e:
            return {"error": f"Failed to extract Office metadata: {str(e)}"}

    def _extract_docx_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract Word document metadata."""
        try:
            from docx import Document

            doc = Document(file_path)
            metadata = {}

            # Core properties
            props = doc.core_properties
            metadata["title"] = props.title or ""
            metadata["author"] = props.author or ""
            metadata["subject"] = props.subject or ""
            metadata["keywords"] = props.keywords or ""
            metadata["category"] = props.category or ""
            metadata["comments"] = props.comments or ""
            metadata["created"] = props.created.isoformat() if props.created else ""
            metadata["modified"] = props.modified.isoformat() if props.modified else ""
            metadata["last_modified_by"] = props.last_modified_by or ""
            metadata["revision"] = props.revision

            # Document structure
            metadata["paragraph_count"] = len(doc.paragraphs)
            metadata["table_count"] = len(doc.tables)

            # Content analysis
            text_content = []
            for para in doc.paragraphs[:10]:  # Sample first 10 paragraphs
                if para.text.strip():
                    text_content.append(para.text.strip())

            metadata["has_content"] = len(text_content) > 0
            metadata["avg_paragraph_length"] = sum(len(p) for p in text_content) // max(len(text_content), 1)

            # Check for images
            image_count = 0
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref.lower():
                    image_count += 1
            metadata["image_count"] = image_count

            return metadata

        except Exception as e:
            return {"error": f"Failed to extract DOCX metadata: {str(e)}"}

    def _extract_xlsx_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract Excel document metadata."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(file_path, read_only=True)
            metadata = {}

            # Core properties
            props = wb.properties
            metadata["title"] = props.title or ""
            metadata["creator"] = props.creator or ""
            metadata["subject"] = props.subject or ""
            metadata["description"] = props.description or ""
            metadata["keywords"] = props.keywords or ""
            metadata["category"] = props.category or ""
            metadata["created"] = props.created.isoformat() if props.created else ""
            metadata["modified"] = props.modified.isoformat() if props.modified else ""
            metadata["last_modified_by"] = props.lastModifiedBy or ""

            # Workbook structure
            metadata["sheet_count"] = len(wb.sheetnames)
            metadata["sheet_names"] = wb.sheetnames

            # Analyze active sheet
            ws = wb.active
            metadata["active_sheet"] = ws.title

            # Count data
            max_row = ws.max_row
            max_col = ws.max_column
            metadata["max_row"] = max_row
            metadata["max_column"] = max_col

            # Sample data analysis
            has_formulas = False
            filled_cells = 0

            for row in ws.iter_rows(max_row=min(100, max_row), max_col=min(20, max_col)):
                for cell in row:
                    if cell.value is not None:
                        filled_cells += 1
                        if isinstance(cell.value, str) and cell.value.startswith('='):
                            has_formulas = True

            metadata["has_formulas"] = has_formulas
            metadata["filled_cells_sample"] = filled_cells
            metadata["data_density"] = filled_cells / (min(100, max_row) * min(20, max_col)) if max_row > 0 and max_col > 0 else 0

            return metadata

        except Exception as e:
            return {"error": f"Failed to extract XLSX metadata: {str(e)}"}

    def _extract_pptx_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract PowerPoint document metadata."""
        try:
            from pptx import Presentation

            pres = Presentation(file_path)
            metadata = {}

            # Core properties
            props = pres.core_properties
            metadata["title"] = props.title or ""
            metadata["author"] = props.author or ""
            metadata["subject"] = props.subject or ""
            metadata["keywords"] = props.keywords or ""
            metadata["category"] = props.category or ""
            metadata["comments"] = props.comments or ""
            metadata["created"] = props.created.isoformat() if props.created else ""
            metadata["modified"] = props.modified.isoformat() if props.modified else ""
            metadata["last_modified_by"] = props.last_modified_by or ""

            # Presentation structure
            metadata["slide_count"] = len(pres.slides)

            # Analyze content
            text_slides = 0
            image_slides = 0
            total_text_length = 0

            for slide in pres.slides[:10]:  # Sample first 10 slides
                has_text = False
                has_image = False

                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        has_text = True
                        total_text_length += len(shape.text)
                    elif shape.shape_type == 13:  # Picture
                        has_image = True

                if has_text:
                    text_slides += 1
                if has_image:
                    image_slides += 1

            metadata["text_slides_sample"] = text_slides
            metadata["image_slides_sample"] = image_slides
            metadata["avg_text_per_slide"] = total_text_length // max(text_slides, 1)

            return metadata

        except Exception as e:
            return {"error": f"Failed to extract PPTX metadata: {str(e)}"}

    def extract_text_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from text files."""
        try:
            metadata = {}

            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            content = None
            used_encoding = None

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue

            if content is None:
                return {"error": "Could not decode file with any common encoding"}

            metadata["encoding"] = used_encoding
            metadata["character_count"] = len(content)
            metadata["line_count"] = content.count('\n') + 1
            metadata["word_count"] = len(content.split())

            # Analyze content structure
            lines = content.split('\n')
            metadata["empty_lines"] = sum(1 for line in lines if not line.strip())
            metadata["avg_line_length"] = sum(len(line) for line in lines) / len(lines) if lines else 0

            # Character set analysis
            unique_chars = set(content)
            metadata["unique_character_count"] = len(unique_chars)

            # Check for special patterns
            has_tabs = '\t' in content
            has_commas = ',' in content and lines
            metadata["has_tab_separation"] = has_tabs
            metadata["has_comma_separation"] = has_commas

            # Detect potential structure
            if has_commas or has_tabs:
                delimiter = '\t' if has_tabs else ','
                # Check if it looks structured
                first_lines = lines[:10]
                column_counts = [len(line.split(delimiter)) for line in first_lines if line.strip()]
                if column_counts and len(set(column_counts)) <= 2:  # Consistent column count
                    metadata["appears_structured"] = True
                    metadata["estimated_columns"] = max(column_counts) if column_counts else 0
                else:
                    metadata["appears_structured"] = False

            return metadata

        except Exception as e:
            return {"error": f"Failed to extract text metadata: {str(e)}"}

    def extract_web_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from web documents."""
        try:
            metadata = {}

            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

            metadata["character_count"] = len(content)
            metadata["line_count"] = content.count('\n') + 1

            # HTML/XML specific analysis
            if content.strip().lower().startswith('<!doctype html') or '<html' in content.lower():
                metadata["document_type"] = "html"

                # Extract HTML metadata
                import re

                # Title
                title_match = re.search(r'<title[^>]*>(.*?)</title>', content, re.IGNORECASE | re.DOTALL)
                metadata["title"] = title_match.group(1).strip() if title_match else ""

                # Meta tags
                meta_matches = re.findall(r'<meta\s+([^>]+)>', content, re.IGNORECASE)
                meta_info = {}
                for meta in meta_matches:
                    if 'name=' in meta and 'content=' in meta:
                        name_match = re.search(r'name=["\']([^"\']+)["\']', meta, re.IGNORECASE)
                        content_match = re.search(r'content=["\']([^"\']+)["\']', meta, re.IGNORECASE)
                        if name_match and content_match:
                            meta_info[name_match.group(1)] = content_match.group(1)

                metadata["meta_tags"] = meta_info

                # Count HTML elements
                metadata["tag_count"] = len(re.findall(r'<[^>]+>', content))
                metadata["has_scripts"] = '<script' in content.lower()
                metadata["has_styles"] = '<style' in content.lower() or 'stylesheet' in content.lower()
                metadata["has_forms"] = '<form' in content.lower()
                metadata["has_tables"] = '<table' in content.lower()

            elif content.strip().startswith('<?xml') or content.strip().startswith('<'):
                metadata["document_type"] = "xml"

                # XML declaration
                xml_decl_match = re.search(r'<\?xml\s+([^>]+)\?>', content)
                if xml_decl_match:
                    metadata["xml_declaration"] = xml_decl_match.group(1)

                # Root element
                root_match = re.search(r'<([a-zA-Z][^>\s]*)', content.lstrip())
                if root_match:
                    metadata["root_element"] = root_match.group(1)

                # Count elements
                metadata["element_count"] = len(re.findall(r'<[^!?/][^>]*>', content))

            else:
                metadata["document_type"] = "unknown_web"

            return metadata

        except Exception as e:
            return {"error": f"Failed to extract web metadata: {str(e)}"}