"""
Office document extractor for DOCX, XLSX, PPTX files.
"""

import logging
from pathlib import Path
from typing import Dict, Any

try:
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
except ImportError as e:
    logging.error(f"Missing required dependencies: {e}")
    raise


def extract_docx(file_path: str) -> Dict[str, Any]:
    """Extract text from DOCX file."""
    try:
        doc = Document(file_path)
        text_parts = []

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text.strip())

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(' | '.join(row_text))

        extracted_text = '\n'.join(text_parts)

        # Process images with Google Vision OCR
        try:
            from universal_image_ocr import UniversalImageOCR
            ocr_processor = UniversalImageOCR(use_google_vision=True)
            ocr_result = ocr_processor.extract_and_process_images(file_path, 'docx')
            
            if ocr_result.get('success') and ocr_result.get('ocr_text'):
                extracted_text += '\n\n=== TEXTO DE IMAGENS (OCR) ===\n\n' + ocr_result['ocr_text']
        except Exception as e:
            # Falha silenciosa - OCR é opcional
            logging.debug(f"DOCX OCR processing failed: {e}")

        return {
            "success": True,
            "extracted_text": extracted_text,
            "error": None
        }

    except Exception as e:
        return {
            "success": False,
            "extracted_text": "",
            "error": f"Error extracting DOCX: {str(e)}"
        }


def extract_xlsx(file_path: str) -> Dict[str, Any]:
    """Extract text from XLSX file."""
    try:
        workbook = load_workbook(file_path, data_only=True)
        text_parts = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"=== Sheet: {sheet_name} ===")

            for row in sheet.iter_rows():
                row_data = []
                for cell in row:
                    if cell.value is not None:
                        row_data.append(str(cell.value))
                if row_data:
                    text_parts.append(' | '.join(row_data))

        extracted_text = '\n'.join(text_parts)

        # Process images with Google Vision OCR
        try:
            from universal_image_ocr import UniversalImageOCR
            ocr_processor = UniversalImageOCR(use_google_vision=True)
            ocr_result = ocr_processor.extract_and_process_images(file_path, 'xlsx')
            
            if ocr_result.get('success') and ocr_result.get('ocr_text'):
                extracted_text += '\n\n=== TEXTO DE IMAGENS (OCR) ===\n\n' + ocr_result['ocr_text']
        except Exception as e:
            # Falha silenciosa - OCR é opcional
            logging.debug(f"XLSX OCR processing failed: {e}")

        return {
            "success": True,
            "extracted_text": extracted_text,
            "error": None
        }

    except Exception as e:
        return {
            "success": False,
            "extracted_text": "",
            "error": f"Error extracting XLSX: {str(e)}"
        }


def extract_pptx(file_path: str) -> Dict[str, Any]:
    """Extract text from PPTX file."""
    try:
        presentation = Presentation(file_path)
        text_parts = []

        for i, slide in enumerate(presentation.slides, 1):
            text_parts.append(f"=== Slide {i} ===")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

        extracted_text = '\n'.join(text_parts)

        # Process images with Google Vision OCR
        try:
            from universal_image_ocr import UniversalImageOCR
            ocr_processor = UniversalImageOCR(use_google_vision=True)
            ocr_result = ocr_processor.extract_and_process_images(file_path, 'pptx')
            
            if ocr_result.get('success') and ocr_result.get('ocr_text'):
                extracted_text += '\n\n=== TEXTO DE IMAGENS (OCR) ===\n\n' + ocr_result['ocr_text']
        except Exception as e:
            # Falha silenciosa - OCR é opcional
            logging.debug(f"PPTX OCR processing failed: {e}")

        return {
            "success": True,
            "extracted_text": extracted_text,
            "error": None
        }

    except Exception as e:
        return {
            "success": False,
            "extracted_text": "",
            "error": f"Error extracting PPTX: {str(e)}"
        }


def extract_office_document(file_path: str) -> Dict[str, Any]:
    """Main function to extract text from office documents."""
    path = Path(file_path)

    if not path.exists():
        return {
            "success": False,
            "extracted_text": "",
            "error": f"File not found: {file_path}"
        }

    extension = path.suffix.lower()

    if extension == '.docx':
        return extract_docx(file_path)
    elif extension == '.xlsx':
        return extract_xlsx(file_path)
    elif extension == '.pptx':
        return extract_pptx(file_path)
    else:
        return {
            "success": False,
            "extracted_text": "",
            "error": f"Unsupported office file format: {extension}"
        }


if __name__ == "__main__":
    import sys
    if len(sys.argv) != 2:
        print("Usage: python office_extractor.py <file_path>")
        sys.exit(1)

    result = extract_office_document(sys.argv[1])
    print(f"Success: {result['success']}")
    if result['success']:
        print(f"Text length: {len(result['extracted_text'])}")
        print(f"First 200 chars: {result['extracted_text'][:200]}...")
    else:
        print(f"Error: {result['error']}")