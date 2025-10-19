#!/usr/bin/env python3
"""
PowerPoint Enhanced Extractor
Extracts text with intelligent chunking by slide + speaker notes.
"""

import sys
import os
import json
import shutil
import tempfile
from pathlib import Path
from typing import Dict, List, Any

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def extract_pptx_enhanced(file_path: str) -> Dict[str, Any]:
    """
    Extract PowerPoint with slide-by-slide chunking.
    
    Returns:
        {
            "success": bool,
            "text": str,
            "slides": [
                {
                    "slide_number": int,
                    "title": str,
                    "content": str,
                    "notes": str,
                    "has_images": bool,
                    "has_tables": bool
                }
            ],
            "chunking_hints": {
                "chunk_by": "slide",
                "preserve_slide_numbers": True
            }
        }
    """
    if not PPTX_AVAILABLE:
        return {
            "success": False,
            "error": "python-pptx not available",
            "text": "",
            "slides": []
        }
    
    temp_link = None
    try:
        path = Path(file_path)
        if not path.exists():
            return {
                "success": False,
                "error": "File not found",
                "text": "",
                "slides": []
            }
        
        # Handle files without extension
        if not file_path.endswith(('.pptx', '.ppt')):
            temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
            temp_file.close()
            shutil.copy2(file_path, temp_file.name)
            file_to_read = temp_file.name
            temp_link = file_to_read
        else:
            file_to_read = file_path
        
        prs = Presentation(file_to_read)
        
        slides_data = []
        text_parts = []
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_info = {
                "slide_number": slide_idx,
                "title": "",
                "content": "",
                "notes": "",
                "has_images": False,
                "has_tables": False
            }
            
            # Extract text from all shapes
            slide_text_parts = []
            
            for shape in slide.shapes:
                # Check if has text
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    
                    # First text is usually the title
                    if not slide_info["title"] and len(text) < 100:
                        slide_info["title"] = text
                    else:
                        slide_text_parts.append(text)
                
                # Check for images
                if shape.shape_type == 13:  # Picture
                    slide_info["has_images"] = True
                
                # Check for tables
                if shape.shape_type == 19:  # Table
                    slide_info["has_tables"] = True
                    # Try to extract table data
                    try:
                        if hasattr(shape, 'table'):
                            table_text = extract_table_from_shape(shape.table)
                            if table_text:
                                slide_text_parts.append(table_text)
                    except:
                        pass
            
            slide_info["content"] = '\n'.join(slide_text_parts)
            
            # Extract speaker notes
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text and notes_text.strip():
                        slide_info["notes"] = notes_text.strip()
                except:
                    pass
            
            # Process images in this slide with Google Vision OCR
            if slide_info["has_images"]:
                try:
                    from universal_image_ocr import UniversalImageOCR
                    ocr_processor = UniversalImageOCR(use_google_vision=True)
                    # Extract images from this specific slide
                    slide_images = ocr_processor._extract_pptx_slide_images(file_path, slide_idx - 1)
                    if slide_images:
                        slide_ocr_text = ""
                        for img_path in slide_images:
                            try:
                                result = ocr_processor.advanced_ocr.process_image(str(img_path))
                                if result.get('success') and result.get('text'):
                                    slide_ocr_text += result['text'] + '\n'
                            except:
                                continue
                        if slide_ocr_text:
                            slide_info["content"] += '\n\n=== TEXTO DE IMAGENS (OCR) ===\n\n' + slide_ocr_text
                except Exception as e:
                    # Falha silenciosa - OCR é opcional
                    pass
            
            slides_data.append(slide_info)
            
            # Build text representation for this slide
            slide_text = f"=== SLIDE {slide_idx}"
            if slide_info["title"]:
                slide_text += f": {slide_info['title']}"
            slide_text += " ==="
            
            text_parts.append(slide_text)
            
            if slide_info["content"]:
                text_parts.append(slide_info["content"])
            
            if slide_info["notes"]:
                text_parts.append(f"Notas: {slide_info['notes']}")
            
            text_parts.append("")  # Empty line between slides
        
        return {
            "success": True,
            "text": '\n'.join(text_parts),
            "slides": slides_data,
            "metadata": {
                "total_slides": len(slides_data),
                "slides_with_notes": sum(1 for s in slides_data if s['notes']),
                "slides_with_images": sum(1 for s in slides_data if s['has_images']),
                "slides_with_tables": sum(1 for s in slides_data if s['has_tables'])
            },
            "chunking_hints": {
                "chunk_by": "slide",
                "preserve_slide_numbers": True
            }
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "text": "",
            "slides": []
        }
    finally:
        # Clean up temp file
        if temp_link and os.path.exists(temp_link):
            try:
                os.remove(temp_link)
            except:
                pass


def extract_table_from_shape(table) -> str:
    """Extract text from PowerPoint table."""
    try:
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cells.append(cell.text.strip() if cell.text else '')
            rows.append(' | '.join(cells))
        
        if rows:
            result = "=== TABELA ===\n"
            result += rows[0] + "\n"  # Header
            result += "-" * 40 + "\n"
            result += '\n'.join(rows[1:])  # Data rows
            return result
        return ""
    except:
        return ""


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"success": False, "error": "No file provided"}))
        sys.exit(1)
    
    file_path = sys.argv[1]
    result = extract_pptx_enhanced(file_path)
    print(json.dumps(result, ensure_ascii=False, indent=2))



