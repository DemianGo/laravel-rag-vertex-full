#!/usr/bin/env python3
"""
Generic Large File Generator for Testing
Generates realistic large files in various formats (PDF, DOCX, XLSX, PPTX, TXT, CSV, HTML, XML, RTF)
No hardcoded test-specific content - all generic and reusable
"""

import os
import sys
from datetime import datetime

def generate_pdf(filepath: str, pages: int = 1000):
    """Generate large PDF file"""
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
        
        c = canvas.Canvas(filepath, pagesize=letter)
        
        for page_num in range(1, pages + 1):
            # Title
            c.setFont("Helvetica-Bold", 16)
            c.drawString(1*inch, 10*inch, f"Test Document - Page {page_num}")
            
            # Body content (generic text)
            c.setFont("Helvetica", 10)
            y_position = 9.5 * inch
            
            for para in range(20):  # 20 paragraphs per page
                text = f"Lorem ipsum dolor sit amet, consectetur adipiscing elit. Paragraph {para+1} on page {page_num}. " * 3
                
                # Wrap text
                words = text.split()
                line = ""
                for word in words:
                    test_line = line + word + " "
                    if c.stringWidth(test_line, "Helvetica", 10) < 6.5 * inch:
                        line = test_line
                    else:
                        c.drawString(1*inch, y_position, line)
                        y_position -= 0.15 * inch
                        line = word + " "
                        
                        if y_position < 1 * inch:
                            break
                
                if line and y_position > 1 * inch:
                    c.drawString(1*inch, y_position, line)
                    y_position -= 0.2 * inch
                
                if y_position < 1 * inch:
                    break
            
            # Footer
            c.setFont("Helvetica", 8)
            c.drawString(1*inch, 0.5*inch, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')} | Page {page_num}/{pages}")
            
            c.showPage()
        
        c.save()
        print(f"✅ PDF criado: {filepath} ({pages} páginas, {os.path.getsize(filepath) / (1024*1024):.2f}MB)")
        return True
        
    except ImportError:
        print("❌ reportlab não instalado. Instale com: pip3 install reportlab")
        return False
    except Exception as e:
        print(f"❌ Erro ao criar PDF: {e}")
        return False

def generate_docx(filepath: str, pages: int = 2000):
    """Generate large DOCX file"""
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        
        doc = Document()
        
        # Add title
        title = doc.add_heading('Large Test Document', 0)
        
        # Estimate: 30 paragraphs per page
        total_paragraphs = pages * 30
        
        for i in range(total_paragraphs):
            if i % 30 == 0:
                # Add section heading every "page"
                page_num = i // 30 + 1
                doc.add_heading(f'Section {page_num}', level=1)
            
            # Add paragraph
            p = doc.add_paragraph(
                f"This is paragraph {i+1} of the test document. Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
                f"Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam, quis nostrud "
                f"exercitation ullamco laboris nisi ut aliquip ex ea commodo consequat. "
            )
            
            # Add page break every 30 paragraphs
            if (i + 1) % 30 == 0 and i < total_paragraphs - 1:
                doc.add_page_break()
        
        doc.save(filepath)
        print(f"✅ DOCX criado: {filepath} (~{pages} páginas, {os.path.getsize(filepath) / (1024*1024):.2f}MB)")
        return True
        
    except ImportError:
        print("❌ python-docx não instalado. Instale com: pip3 install python-docx")
        return False
    except Exception as e:
        print(f"❌ Erro ao criar DOCX: {e}")
        return False

def generate_xlsx(filepath: str, rows: int = 10000):
    """Generate large XLSX file"""
    try:
        from openpyxl import Workbook
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        
        # Headers
        headers = ['ID', 'Name', 'Value', 'Date', 'Description', 'Category', 'Status', 'Notes']
        ws.append(headers)
        
        # Data rows
        for i in range(1, rows + 1):
            row = [
                i,
                f"Item_{i:06d}",
                round(i * 1.5, 2),
                datetime.now().strftime('%Y-%m-%d'),
                f"Description for row {i} with additional text to increase size",
                f"Category_{i % 10}",
                "Active" if i % 2 == 0 else "Inactive",
                f"Notes for item {i}: Lorem ipsum dolor sit amet, consectetur adipiscing elit."
            ]
            ws.append(row)
        
        wb.save(filepath)
        print(f"✅ XLSX criado: {filepath} ({rows} linhas, {os.path.getsize(filepath) / (1024*1024):.2f}MB)")
        return True
        
    except ImportError:
        print("❌ openpyxl não instalado. Instale com: pip3 install openpyxl")
        return False
    except Exception as e:
        print(f"❌ Erro ao criar XLSX: {e}")
        return False

def generate_pptx(filepath: str, slides: int = 500):
    """Generate large PPTX file"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation()
        
        for i in range(1, slides + 1):
            # Add slide with title and content layout
            slide_layout = prs.slide_layouts[1]  # Title and content
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = f"Slide {i}: Test Content"
            
            # Content
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = f"This is slide {i} of the test presentation."
            
            # Add bullet points
            for j in range(5):
                p = tf.add_paragraph()
                p.text = f"Bullet point {j+1}: Lorem ipsum dolor sit amet, consectetur adipiscing elit. Additional content to increase size."
                p.level = 1
        
        prs.save(filepath)
        print(f"✅ PPTX criado: {filepath} ({slides} slides, {os.path.getsize(filepath) / (1024*1024):.2f}MB)")
        return True
        
    except ImportError:
        print("❌ python-pptx não instalado. Instale com: pip3 install python-pptx")
        return False
    except Exception as e:
        print(f"❌ Erro ao criar PPTX: {e}")
        return False

def generate_txt(filepath: str, lines: int = 50000):
    """Generate large TXT file"""
    try:
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(f"Large Test Text File\n")
            f.write(f"Generated: {datetime.now()}\n")
            f.write("=" * 80 + "\n\n")
            
            for i in range(1, lines + 1):
                f.write(f"Line {i}: Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
                       f"Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. "
                       f"Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris.\n")
                
                # Add section header every 100 lines
                if i % 100 == 0:
                    f.write("\n" + "=" * 80 + "\n")
                    f.write(f"SECTION {i//100}\n")
                    f.write("=" * 80 + "\n\n")
        
        print(f"✅ TXT criado: {filepath} ({lines} linhas, {os.path.getsize(filepath) / (1024*1024):.2f}MB)")
        return True
        
    except Exception as e:
        print(f"❌ Erro ao criar TXT: {e}")
        return False

def generate_csv(filepath: str, rows: int = 50000):
    """Generate large CSV file"""
    try:
        import csv
        
        with open(filepath, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            
            # Header
            writer.writerow(['ID', 'Name', 'Email', 'Value', 'Date', 'Description', 'Category', 'Status'])
            
            # Data rows
            for i in range(1, rows + 1):
                writer.writerow([
                    i,
                    f"User_{i:06d}",
                    f"user{i}@example.com",
                    round(i * 1.23, 2),
                    datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
                    f"Description for record {i} with additional text to increase file size",
                    f"Cat_{i % 20}",
                    "Active" if i % 3 == 0 else "Pending"
                ])
        
        print(f"✅ CSV criado: {filepath} ({rows} linhas, {os.path.getsize(filepath) / (1024*1024):.2f}MB)")
        return True
        
    except Exception as e:
        print(f"❌ Erro ao criar CSV: {e}")
        return False

def generate_html(filepath: str, sections: int = 1000):
    """Generate large HTML file"""
    try:
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write('<!DOCTYPE html>\n<html lang="pt-BR">\n<head>\n')
            f.write('    <meta charset="UTF-8">\n')
            f.write('    <title>Large Test HTML Document</title>\n')
            f.write('    <style>body { font-family: Arial, sans-serif; margin: 40px; }</style>\n')
            f.write('</head>\n<body>\n')
            f.write(f'    <h1>Large Test HTML Document</h1>\n')
            f.write(f'    <p>Generated: {datetime.now()}</p>\n')
            
            for i in range(1, sections + 1):
                f.write(f'    <h2>Section {i}</h2>\n')
                f.write(f'    <p>This is section {i} of the test document. Lorem ipsum dolor sit amet, consectetur adipiscing elit. ')
                f.write('Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. Ut enim ad minim veniam, quis nostrud ')
                f.write('exercitation ullamco laboris nisi ut aliquip ex ea commodo consequat.</p>\n')
                
                # Add list
                f.write('    <ul>\n')
                for j in range(5):
                    f.write(f'        <li>Item {j+1}: Additional content to increase file size</li>\n')
                f.write('    </ul>\n')
            
            f.write('</body>\n</html>')
        
        print(f"✅ HTML criado: {filepath} ({sections} seções, {os.path.getsize(filepath) / (1024*1024):.2f}MB)")
        return True
        
    except Exception as e:
        print(f"❌ Erro ao criar HTML: {e}")
        return False

def generate_xml(filepath: str, records: int = 10000):
    """Generate large XML file"""
    try:
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write('<?xml version="1.0" encoding="UTF-8"?>\n')
            f.write('<data>\n')
            f.write(f'    <metadata>\n')
            f.write(f'        <generated>{datetime.now()}</generated>\n')
            f.write(f'        <records>{records}</records>\n')
            f.write(f'    </metadata>\n')
            f.write('    <items>\n')
            
            for i in range(1, records + 1):
                f.write(f'        <item id="{i}">\n')
                f.write(f'            <name>Item_{i:06d}</name>\n')
                f.write(f'            <value>{i * 1.5:.2f}</value>\n')
                f.write(f'            <date>{datetime.now().strftime("%Y-%m-%d")}</date>\n')
                f.write(f'            <description>Description for item {i} with additional text to increase size</description>\n')
                f.write(f'            <category>Category_{i % 10}</category>\n')
                f.write(f'            <status>{"active" if i % 2 == 0 else "inactive"}</status>\n')
                f.write(f'        </item>\n')
            
            f.write('    </items>\n')
            f.write('</data>')
        
        print(f"✅ XML criado: {filepath} ({records} registros, {os.path.getsize(filepath) / (1024*1024):.2f}MB)")
        return True
        
    except Exception as e:
        print(f"❌ Erro ao criar XML: {e}")
        return False

def generate_rtf(filepath: str, pages: int = 1000):
    """Generate large RTF file"""
    try:
        with open(filepath, 'w', encoding='utf-8') as f:
            # RTF header
            f.write(r'{\rtf1\ansi\deff0' + '\n')
            f.write(r'{\fonttbl{\f0 Times New Roman;}}' + '\n')
            f.write(r'{\colortbl;\red0\green0\blue0;}' + '\n')
            
            # Content
            for i in range(1, pages * 30 + 1):  # 30 paragraphs per page
                f.write(r'\par ')
                f.write(f'Paragraph {i}: Lorem ipsum dolor sit amet, consectetur adipiscing elit. ')
                f.write('Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua. ')
                f.write('Ut enim ad minim veniam, quis nostrud exercitation ullamco laboris. ')
                
                # Page break every 30 paragraphs
                if i % 30 == 0:
                    f.write(r'\page ')
            
            f.write(r'}')
        
        print(f"✅ RTF criado: {filepath} (~{pages} páginas, {os.path.getsize(filepath) / (1024*1024):.2f}MB)")
        return True
        
    except Exception as e:
        print(f"❌ Erro ao criar RTF: {e}")
        return False

def main():
    output_dir = "/tmp/large_test_files"
    os.makedirs(output_dir, exist_ok=True)
    
    print("═══════════════════════════════════════════════════════════════════════")
    print("🔧 Gerando arquivos de teste gigantes (genéricos)")
    print("═══════════════════════════════════════════════════════════════════════")
    print("")
    
    generators = [
        ("PDF 1000 páginas", lambda: generate_pdf(f"{output_dir}/test_1000pages.pdf", 1000)),
        ("PDF 3000 páginas", lambda: generate_pdf(f"{output_dir}/test_3000pages.pdf", 3000)),
        ("PDF 5000 páginas", lambda: generate_pdf(f"{output_dir}/test_5000pages.pdf", 5000)),
        ("DOCX 2000 páginas", lambda: generate_docx(f"{output_dir}/test_2000pages.docx", 2000)),
        ("XLSX 10000 linhas", lambda: generate_xlsx(f"{output_dir}/test_10000rows.xlsx", 10000)),
        ("PPTX 500 slides", lambda: generate_pptx(f"{output_dir}/test_500slides.pptx", 500)),
        ("TXT 50000 linhas", lambda: generate_txt(f"{output_dir}/test_50000lines.txt", 50000)),
        ("CSV 50000 linhas", lambda: generate_csv(f"{output_dir}/test_50000rows.csv", 50000)),
        ("HTML 1000 seções", lambda: generate_html(f"{output_dir}/test_1000sections.html", 1000)),
        ("XML 10000 registros", lambda: generate_xml(f"{output_dir}/test_10000records.xml", 10000)),
        ("RTF 1000 páginas", lambda: generate_rtf(f"{output_dir}/test_1000pages.rtf", 1000)),
    ]
    
    success_count = 0
    for name, generator in generators:
        print(f"\n📝 Gerando: {name}")
        if generator():
            success_count += 1
    
    print("")
    print("═══════════════════════════════════════════════════════════════════════")
    print(f"✅ Geração concluída: {success_count}/{len(generators)} arquivos criados")
    print("═══════════════════════════════════════════════════════════════════════")
    print(f"\n📁 Arquivos salvos em: {output_dir}")
    print("")

if __name__ == '__main__':
    main()

